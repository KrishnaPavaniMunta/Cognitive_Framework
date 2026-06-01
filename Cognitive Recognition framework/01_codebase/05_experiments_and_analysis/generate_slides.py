"""
Generate a management-friendly PowerPoint status update for the Porto meeting.
Styled to match "Role of Data and AI - WP3.pptx" template.

Usage:  python generate_slides.py
Output: hospital_ai_status.pptx  (workspace root)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Colours — matched to the WP3 template ──────────────────────────────────────
NAVY       = RGBColor(0x1D, 0x1E, 0x37)   # primary title / dark text
ORANGE     = RGBColor(0xFA, 0x64, 0x3F)   # accent (matches template #FA643F)
BLUE_STEEL = RGBColor(0x4F, 0x81, 0xBD)   # theme accent1
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xFA, 0xFA, 0xFC)
MID_GREY   = RGBColor(0x88, 0x88, 0x88)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
AMBER      = RGBColor(0xE6, 0x7E, 0x22)

FONT = "Calibri"

# ── Slide dimensions — 20 × 11.25 in (from WP3 template) ───────────────────────
W = Inches(20.0)
H = Inches(11.25)

HEADER_H  = Inches(1.6)
CONTENT_Y = Inches(1.75)
CONTENT_H = H - CONTENT_Y - Inches(0.35)
MARGIN    = Inches(0.75)
CONTENT_W = W - 2 * MARGIN

WORKSPACE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE  = os.path.join(WORKSPACE, "Role of Data and AI- WP3.pptx")
IMG_DIR   = os.path.join(WORKSPACE, "outputs", "hospital_results")
OUT_FILE  = os.path.join(WORKSPACE, "hospital_ai_status.pptx")


# ── Helpers ────────────────────────────────────────────────────────────────────

def open_template():
    """Open the WP3 template and remove its existing 6 slides."""
    prs = Presentation(TEMPLATE)
    sldIdLst = prs.slides._sldIdLst
    for i in range(len(prs.slides) - 1, -1, -1):
        sldId_elem = sldIdLst[i]
        rId = sldId_elem.get(qn('r:id'))
        prs.slides.part.drop_rel(rId)
        sldIdLst.remove(sldId_elem)
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # Blank layout


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line_color=None, line_pt=1.5):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    fill_shape(shape, color)
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=20, bold=False, color=NAVY,
                 align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = FONT
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_header(slide, title, subtitle=None):
    """Orange left accent bar + bold navy title, matching template style."""
    add_rect(slide, 0, 0, Inches(0.22), HEADER_H, ORANGE)
    add_text_box(slide, title,
                 MARGIN, Inches(0.2), W - MARGIN - Inches(0.3), Inches(0.9),
                 font_size=44, bold=True, color=NAVY)
    add_rect(slide, MARGIN, Inches(1.25), W - 2 * MARGIN, Inches(0.06), ORANGE)
    if subtitle:
        add_text_box(slide, subtitle,
                     MARGIN, Inches(1.35), W - 2 * MARGIN, Inches(0.38),
                     font_size=22, italic=True,
                     color=RGBColor(0x55, 0x55, 0x55))


# ── Individual slide builders ──────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs)
    # Navy left panel
    add_rect(slide, 0, 0, Inches(7.0), H, NAVY)
    add_rect(slide, Inches(7.0), 0, Inches(0.2), H, ORANGE)

    add_text_box(slide,
                 "AI-Powered Hospital\nObject Detection",
                 Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.8),
                 font_size=56, bold=True, color=WHITE)

    add_rect(slide, Inches(0.6), Inches(6.0), Inches(5.0), Inches(0.1), ORANGE)

    add_text_box(slide, "Project Status Update",
                 Inches(0.6), Inches(6.2), Inches(6.0), Inches(0.75),
                 font_size=30, color=RGBColor(0xAA, 0xBB, 0xCC))

    add_text_box(slide, "April 2026  \u00b7  Porto Project Meeting",
                 Inches(0.6), Inches(7.05), Inches(6.0), Inches(0.6),
                 font_size=24, color=RGBColor(0x88, 0x99, 0xAA))

    # Right panel: big stat
    add_text_box(slide, "109",
                 Inches(9.5), Inches(2.2), Inches(8.5), Inches(3.5),
                 font_size=180, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "object categories detected",
                 Inches(9.0), Inches(5.8), Inches(9.5), Inches(0.8),
                 font_size=30, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "across 5 hospital-relevant groups",
                 Inches(9.0), Inches(6.65), Inches(9.5), Inches(0.6),
                 font_size=24, italic=True,
                 color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)


def slide_overview(prs):
    slide = blank_slide(prs)
    add_header(slide, "Project Overview",
               "What we are building and why")

    boxes = [
        ("What",
         "A computer vision system that automatically recognises objects, people, and safety-relevant items in hospital environments using video camera feeds."),
        ("Why",
         "Supports staff situational awareness, safety monitoring, and operational efficiency \u2014 without manual review of footage."),
        ("How",
         "We fine-tune a state-of-the-art YOLO detection model on a custom hospital dataset, using a two-phase training strategy that preserves general knowledge while learning hospital-specific objects."),
    ]

    bw = (CONTENT_W - Inches(0.6)) / 3
    for i, (label, body) in enumerate(boxes):
        bx = MARGIN + i * (bw + Inches(0.3))
        by = CONTENT_Y + Inches(0.15)
        bh = CONTENT_H - Inches(0.15)

        # Drop shadow
        add_rect(slide, bx + Inches(0.1), by + Inches(0.1), bw, bh,
                 RGBColor(0xDD, 0xDD, 0xDD))
        # Card
        add_rect(slide, bx, by, bw, bh, WHITE, NAVY, 1.5)
        # Orange top bar
        add_rect(slide, bx, by, bw, Inches(0.72), ORANGE)
        add_text_box(slide, label, bx, by + Inches(0.1), bw, Inches(0.58),
                     font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, body,
                     bx + Inches(0.3), by + Inches(0.9),
                     bw - Inches(0.6), bh - Inches(1.1),
                     font_size=21, color=RGBColor(0x22, 0x22, 0x22))


def slide_capabilities(prs):
    slide = blank_slide(prs)
    add_header(slide, "Detection Capabilities",
               "109 object categories across 5 groups")

    categories = [
        ("Medical\nEquipment",
         ["Hospital beds\n& stretchers", "IV stands &\ninfusion pumps",
          "Patient monitors", "Surgical lights", "Nasal cannulas"]),
        ("Staff &\nPatients",
         ["Healthcare workers", "Patients",
          "Hair nets,\nmasks & gloves"]),
        ("Safety &\nAccess",
         ["Fire extinguishers", "Exit signs", "Wheelchairs",
          "Spillage detection", "Security cameras"]),
        ("Hospital\nInfrastructure",
         ["Doors & corridors", "Reception counters", "Bathroom labels",
          "Cabinets & benches", "Vending machines"]),
        ("Common /\nBackground",
         ["People & vehicles", "Bags &\neveryday items",
          "80 standard COCO\ncategories retained"]),
    ]

    colors = [NAVY, ORANGE, BLUE_STEEL,
              RGBColor(0x6C, 0x3B, 0x93), RGBColor(0x2E, 0x86, 0x48)]

    bw = (CONTENT_W - Inches(0.8)) / 5
    for i, (cat, items) in enumerate(categories):
        bx = MARGIN + i * (bw + Inches(0.2))
        by = CONTENT_Y
        bh = CONTENT_H

        add_rect(slide, bx, by, bw, Inches(1.0), colors[i])
        add_text_box(slide, cat, bx + Inches(0.05), by + Inches(0.08),
                     bw - Inches(0.1), Inches(0.88),
                     font_size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        card = add_rect(slide, bx, by + Inches(1.0),
                        bw, bh - Inches(1.0), WHITE)
        card.line.color.rgb = colors[i]
        card.line.width = Pt(1.5)

        for j, item in enumerate(items):
            iy = by + Inches(1.15) + j * Inches(1.4)
            add_rect(slide, bx + Inches(0.2), iy + Inches(0.2),
                     Inches(0.25), Inches(0.25), colors[i])
            add_text_box(slide, item,
                         bx + Inches(0.58), iy,
                         bw - Inches(0.72), Inches(1.2),
                         font_size=19, color=RGBColor(0x22, 0x22, 0x22))


def slide_dataset(prs):
    slide = blank_slide(prs)
    add_header(slide, "Dataset Composition",
               "Three complementary data sources, carefully balanced")

    sources = [
        ("COCO 2017\n(Public Dataset)",
         "~5,000 validation images covering 80 everyday object types. Gives the model general visual understanding so hospital-specific training preserves existing knowledge.",
         NAVY),
        ("Custom Hospital Data\n(Internal Collection)",
         "Images from real hospital environments: beds, equipment, staff, and infrastructure. Filtered and curated to remove ambiguous or low-quality annotations.",
         ORANGE),
        ("Roboflow Specialist Data\n(6 targeted classes)",
         "Publicly available datasets for: bags, doors, exit signs, fire extinguishers, spillage, and wheelchairs. Rare classes upsampled up to \u00d73 to prevent class imbalance.",
         BLUE_STEEL),
    ]

    bw = (CONTENT_W - Inches(0.6)) / 3
    for i, (title, body, col) in enumerate(sources):
        bx = MARGIN + i * (bw + Inches(0.3))
        by = CONTENT_Y + Inches(0.1)
        bh = Inches(5.5)

        add_rect(slide, bx + Inches(0.1), by + Inches(0.1), bw, bh,
                 RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(slide, bx, by, bw, Inches(0.95), col)
        add_text_box(slide, title, bx + Inches(0.2), by + Inches(0.1),
                     bw - Inches(0.4), Inches(0.8),
                     font_size=22, bold=True, color=WHITE)

        card = add_rect(slide, bx, by + Inches(0.95), bw, bh - Inches(0.95), WHITE)
        card.line.color.rgb = col
        card.line.width = Pt(1.5)
        add_text_box(slide, body,
                     bx + Inches(0.25), by + Inches(1.12),
                     bw - Inches(0.5), bh - Inches(1.35),
                     font_size=20, color=RGBColor(0x22, 0x22, 0x22))

    # Stats row
    ry = CONTENT_Y + Inches(5.9)
    stats = [
        ("109 \u2192", "Total object categories"),
        ("85 / 15 %", "Train / Validation split"),
        ("\u00d73", "Oversampling on rarest classes"),
        ("3", "Independent data sources merged"),
    ]
    sw = CONTENT_W / 4
    for i, (num, label) in enumerate(stats):
        sx = MARGIN + i * sw
        add_rect(slide, sx, ry, sw - Inches(0.12), Inches(2.2), NAVY)
        add_text_box(slide, num,
                     sx, ry + Inches(0.1), sw - Inches(0.12), Inches(1.05),
                     font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     sx, ry + Inches(1.15), sw - Inches(0.12), Inches(0.9),
                     font_size=19, color=RGBColor(0xAA, 0xCC, 0xEE),
                     align=PP_ALIGN.CENTER)


def slide_training(prs):
    slide = blank_slide(prs)
    add_header(slide, "Training Approach",
               "Two-phase strategy: learn new knowledge without forgetting the old")

    phases = [
        ("Phase 1", "Teach New Objects\n(Head-Only Training)",
         "Only the final detection layers are updated.\n\n"
         "Duration:  30 epochs  (~46 hours)\n\n"
         "Goal: Rapidly teach the model to recognise hospital-specific objects while keeping the core visual understanding intact.",
         "\u2713   Completed", GREEN),
        ("Phase 2", "Refine & Generalise\n(Full Fine-Tuning)",
         "The deeper feature-extraction layers (neck) are also unlocked and trained.\n\n"
         "Duration:  70 epochs  (in progress)\n\n"
         "Goal: Adapt the model's internal representations to hospital environments, further boosting precision and recall.",
         "\U0001f504   In Progress", AMBER),
    ]

    bw = (CONTENT_W - Inches(0.6)) / 2
    for i, (phase, sub, body, status, scol) in enumerate(phases):
        bx = MARGIN + i * (bw + Inches(0.6))
        by = CONTENT_Y + Inches(0.1)
        bh = CONTENT_H - Inches(0.15)

        add_rect(slide, bx + Inches(0.1), by + Inches(0.1), bw, bh,
                 RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(slide, bx, by, bw, Inches(0.8), NAVY)
        add_text_box(slide, phase, bx, by + Inches(0.1), bw, Inches(0.62),
                     font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        card = add_rect(slide, bx, by + Inches(0.8), bw, bh - Inches(0.8), WHITE)
        card.line.color.rgb = NAVY
        card.line.width = Pt(1.5)

        add_text_box(slide, sub,
                     bx + Inches(0.35), by + Inches(1.0),
                     bw - Inches(0.7), Inches(1.2),
                     font_size=26, bold=True, color=NAVY)

        add_text_box(slide, body,
                     bx + Inches(0.35), by + Inches(2.2),
                     bw - Inches(0.7), bh - Inches(3.15),
                     font_size=22, color=RGBColor(0x22, 0x22, 0x22))

        add_rect(slide, bx + Inches(0.3), by + bh - Inches(0.72),
                 bw - Inches(0.6), Inches(0.6), scol)
        add_text_box(slide, status,
                     bx + Inches(0.3), by + bh - Inches(0.72),
                     bw - Inches(0.6), Inches(0.6),
                     font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_performance(prs):
    slide = blank_slide(prs)
    add_header(slide, "Model Performance",
               "Hospital V3 model \u00b7 Phase 1 results (epoch 30 of 30)")

    metrics = [
        ("90.7 %",  "Precision",
         "When the model flags an object,\nit is correct 9 times out of 10."),
        ("85.8 %",  "Recall",
         "The model finds 86 out of every\n100 objects in the scene."),
        ("91.7 %",  "Accuracy Score\n(mAP @ 50 % overlap)",
         "Industry-standard benchmark\nfor overall detection quality."),
        ("+43 %",   "Improvement vs. V2",
         "V3 Phase 1 already beats\nthe fully-trained previous model."),
    ]

    tw = (CONTENT_W - Inches(0.6)) / 4
    for i, (num, label, note) in enumerate(metrics):
        tx = MARGIN + i * (tw + Inches(0.2))
        ty = CONTENT_Y + Inches(0.1)
        th = Inches(4.5)

        add_rect(slide, tx + Inches(0.07), ty + Inches(0.07), tw, th,
                 RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(slide, tx, ty, tw, th, NAVY)

        num_color = ORANGE if i == 3 else WHITE
        add_text_box(slide, num,
                     tx + Inches(0.05), ty + Inches(0.25),
                     tw - Inches(0.1), Inches(1.65),
                     font_size=64, bold=True, color=num_color, align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     tx + Inches(0.12), ty + Inches(1.95),
                     tw - Inches(0.24), Inches(0.9),
                     font_size=21, bold=True,
                     color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
        add_text_box(slide, note,
                     tx + Inches(0.12), ty + Inches(2.9),
                     tw - Inches(0.24), Inches(1.4),
                     font_size=19, color=RGBColor(0xCC, 0xDD, 0xEE),
                     align=PP_ALIGN.CENTER)

    # Comparison bar chart
    cy = CONTENT_Y + Inches(4.85)
    ch = H - cy - Inches(0.3)
    add_rect(slide, MARGIN, cy, CONTENT_W, ch, OFF_WHITE,
             RGBColor(0xCC, 0xCC, 0xCC), 1)

    add_text_box(slide,
                 "Head-to-head comparison  (mAP50-95, higher = better)",
                 MARGIN + Inches(0.3), cy + Inches(0.18),
                 CONTENT_W - Inches(0.6), Inches(0.5),
                 font_size=22, bold=True, color=NAVY)

    bars = [
        ("V2 \u2014 fully trained, 70 epochs", 0.4731, MID_GREY),
        ("V3 \u2014 Phase 1 only, 30 epochs  (still improving!)", 0.6753, ORANGE),
    ]
    label_w = Inches(6.5)
    max_bar = CONTENT_W - label_w - Inches(2.0)
    for j, (lbl, val, col) in enumerate(bars):
        by2 = cy + Inches(0.85) + j * Inches(1.3)
        add_text_box(slide, lbl,
                     MARGIN + Inches(0.3), by2,
                     label_w, Inches(0.6),
                     font_size=21, color=RGBColor(0x22, 0x22, 0x22))
        bar_w = max_bar * val
        add_rect(slide,
                 MARGIN + label_w + Inches(0.2), by2 + Inches(0.06),
                 bar_w, Inches(0.52), col)
        add_text_box(slide,
                     f"{val:.1%}",
                     MARGIN + label_w + Inches(0.2) + bar_w + Inches(0.15),
                     by2, Inches(1.1), Inches(0.6),
                     font_size=21, bold=True, color=col)


def slide_progress(prs):
    slide = blank_slide(prs)
    add_header(slide, "Project Progress",
               "Milestone tracker \u2014 April 2026")

    milestones = [
        ("\u2713", "Done",        GREEN,    "Data Pipeline",
         "Three datasets downloaded, cleaned, merged and split into training / validation sets."),
        ("\u2713", "Done",        GREEN,    "Hospital V2 Model (106 classes) \u2014 Phases 1 & 2",
         "Both training phases completed. Model validated and weights saved."),
        ("\u2713", "Done",        GREEN,    "Hospital V3 Dataset (109 classes) prepared",
         "Three new specialist classes added. Dataset rebalanced with oversampling."),
        ("\u2713", "Done",        GREEN,    "V3 Phase 1 \u2014 Head Training (30 epochs)",
         "mAP50 = 91.7 %  \u00b7  mAP50-95 = 67.5 %  \u00b7  already 43 % stronger than the full V2 model."),
        ("\U0001f504", "In Progress", AMBER, "V3 Phase 2 \u2014 Full Fine-Tuning (70 epochs)",
         "Training script ready. Neck + head layers unlocked for deeper adaptation."),
        ("\u25cb", "Pending",    MID_GREY,  "Final Validation & Benchmarking",
         "End-to-end evaluation on held-out hospital footage once Phase 2 completes."),
        ("\u25cb", "Pending",    MID_GREY,  "Deployment & Integration Planning",
         "Model packaging for edge / server deployment; integration with camera feeds."),
    ]

    row_h = (CONTENT_H - Inches(0.1)) / len(milestones)
    for i, (sym, slabel, scol, title, detail) in enumerate(milestones):
        ry = CONTENT_Y + Inches(0.05) + i * row_h
        bg_col = WHITE if i % 2 == 0 else OFF_WHITE
        add_rect(slide, MARGIN, ry, CONTENT_W, row_h - Inches(0.06), bg_col)

        badge_w = Inches(1.6)
        add_rect(slide, MARGIN, ry + Inches(0.08),
                 badge_w, row_h - Inches(0.18), scol)
        add_text_box(slide, slabel,
                     MARGIN, ry + Inches(0.1), badge_w, Inches(0.42),
                     font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, sym,
                     MARGIN, ry + Inches(0.48), badge_w, Inches(0.5),
                     font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text_box(slide, title,
                     MARGIN + Inches(1.82), ry + Inches(0.1),
                     Inches(8.5), Inches(0.55),
                     font_size=22, bold=True, color=NAVY)
        add_text_box(slide, detail,
                     MARGIN + Inches(1.82), ry + Inches(0.6),
                     CONTENT_W - Inches(1.95), Inches(0.45),
                     font_size=18, color=RGBColor(0x33, 0x33, 0x33))


def slide_demo(prs):
    slide = blank_slide(prs)
    add_header(slide, "Sample Detections",
               "Model output on real hospital images \u2014 bounding boxes drawn automatically")

    image_files = [
        ("doctors-hospital-corridor-nurse-pushing-gurney-stretcher-bed_annotated.jpg",
         "Corridor: stretcher, staff, door"),
        ("belova59-laboratory-3827738_1920_annotated.jpg",
         "Laboratory: equipment & personnel"),
        ("pexels-photo-4421486_annotated.jpeg",
         "Patient care: bed, monitor, IV stand"),
    ]

    iw = (CONTENT_W - Inches(0.6)) / 3
    ih = CONTENT_H - Inches(0.75)

    for i, (fname, caption) in enumerate(image_files):
        ix = MARGIN + i * (iw + Inches(0.3))
        iy = CONTENT_Y + Inches(0.1)
        fpath = os.path.join(IMG_DIR, fname)

        add_rect(slide, ix + Inches(0.08), iy + Inches(0.08), iw, ih + Inches(0.6),
                 RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(slide, ix, iy, iw, ih + Inches(0.6), WHITE, NAVY, 1.5)

        if os.path.exists(fpath):
            try:
                slide.shapes.add_picture(fpath,
                                         ix + Inches(0.06), iy + Inches(0.06),
                                         iw - Inches(0.12), ih)
            except Exception:
                add_text_box(slide, "[Preview unavailable]",
                             ix, iy + ih / 2, iw, Inches(0.55),
                             font_size=18, color=MID_GREY, align=PP_ALIGN.CENTER)
        else:
            add_text_box(slide, f"[{fname}]",
                         ix + Inches(0.1), iy + Inches(0.1),
                         iw - Inches(0.2), ih,
                         font_size=15, color=MID_GREY)

        # Caption strip
        add_rect(slide, ix, iy + ih, iw, Inches(0.6), NAVY)
        add_text_box(slide, caption,
                     ix + Inches(0.1), iy + ih + Inches(0.1),
                     iw - Inches(0.2), Inches(0.45),
                     font_size=20, bold=False, color=WHITE, align=PP_ALIGN.CENTER)


def slide_next_steps(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, Inches(5.6), H, NAVY)
    add_rect(slide, Inches(5.6), 0, Inches(0.2), H, ORANGE)

    add_text_box(slide, "Next\nSteps",
                 Inches(0.5), Inches(2.4), Inches(4.8), Inches(3.2),
                 font_size=72, bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(5.75), Inches(4.0), Inches(0.1), ORANGE)
    add_text_box(slide, "Priorities after\nPorto meeting",
                 Inches(0.5), Inches(5.95), Inches(4.8), Inches(1.0),
                 font_size=24, italic=True, color=RGBColor(0x88, 0x99, 0xAA))

    steps = [
        ("1", "Launch V3 Phase 2 Training",
         "Start the 70-epoch full fine-tuning run (neck + head layers). Monitor metrics at regular checkpoints."),
        ("2", "Final Validation on Hospital Footage",
         "Evaluate the model on held-out hospital video clips. Measure per-class precision and recall. Flag underperforming categories."),
        ("3", "Deployment & Integration Planning",
         "Assess deployment targets (edge device vs. server). Define API / streaming integration. Package model in ONNX / TensorRT for production."),
    ]

    step_h = (H - Inches(0.5)) / 3
    for i, (num, title, body) in enumerate(steps):
        sy = Inches(0.25) + i * step_h

        add_rect(slide, Inches(6.1), sy + Inches(0.25),
                 Inches(0.88), Inches(0.88), ORANGE)
        add_text_box(slide, num,
                     Inches(6.1), sy + Inches(0.25),
                     Inches(0.88), Inches(0.88),
                     font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text_box(slide, title,
                     Inches(7.2), sy + Inches(0.2),
                     W - Inches(7.6), Inches(0.7),
                     font_size=26, bold=True, color=NAVY)
        add_text_box(slide, body,
                     Inches(7.2), sy + Inches(0.88),
                     W - Inches(7.6), step_h - Inches(1.1),
                     font_size=22, color=RGBColor(0x33, 0x33, 0x33))

        if i < 2:
            add_rect(slide, Inches(5.95), sy + step_h - Inches(0.05),
                     W - Inches(6.2), Inches(0.04),
                     RGBColor(0xCC, 0xCC, 0xCC))


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = open_template()

    print("Building slides...")
    slide_title(prs)         ; print("  1/9  Title")
    slide_overview(prs)      ; print("  2/9  Overview")
    slide_capabilities(prs)  ; print("  3/9  Capabilities")
    slide_dataset(prs)       ; print("  4/9  Dataset")
    slide_training(prs)      ; print("  5/9  Training approach")
    slide_performance(prs)   ; print("  6/9  Performance")
    slide_progress(prs)      ; print("  7/9  Progress tracker")
    slide_demo(prs)          ; print("  8/9  Sample detections")
    slide_next_steps(prs)    ; print("  9/9  Next steps")

    prs.save(OUT_FILE)
    print(f"\nSaved: {OUT_FILE}")


if __name__ == "__main__":
    main()
