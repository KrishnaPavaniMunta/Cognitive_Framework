"""
generate_quarterly_review_slides.py
====================================
Generates: HospitalGuard_QuarterlyReview_May2026.pptx

7-slide quarterly review deck for non-technical project manager audience.
Styled to match "Role of Data and AI - WP3.pptx" (navy / orange / Calibri).

Slides
  1. Title & Executive Summary
  2. The Vision — Health & Safety in Hospitals
  3. The Cognitive Engine & Context Verification
  4. Advanced Tracking — The AI's Memory
  5. Spatial 3D Mapping & Smart Filtering
  6. Off-Board Processing Architecture
  7. Next Phase & Future Deployment

Usage:
    python generate_quarterly_review_slides.py

Output:
    05_documents_and_presentations/HospitalGuard_QuarterlyReview_May2026.pptx
"""

from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Colour palette (matches WP3 template) ─────────────────────────────────────
NAVY        = RGBColor(0x1D, 0x1E, 0x37)
ORANGE      = RGBColor(0xFA, 0x64, 0x3F)
BLUE_STEEL  = RGBColor(0x4F, 0x81, 0xBD)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF5, 0xF6, 0xFA)
MID_GREY    = RGBColor(0x88, 0x88, 0x88)
SHADOW      = RGBColor(0xCC, 0xCC, 0xCC)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
AMBER       = RGBColor(0xE6, 0x7E, 0x22)
PURPLE      = RGBColor(0x6C, 0x3B, 0x93)
TEAL        = RGBColor(0x1A, 0x9E, 0x8C)
DARK_NAVY2  = RGBColor(0x2D, 0x2F, 0x52)
NEAR_BLACK  = RGBColor(0x12, 0x14, 0x28)
PH_BORDER   = RGBColor(0x4F, 0x81, 0xBD)
TEXT_DARK   = RGBColor(0x22, 0x22, 0x22)
TEXT_LIGHT  = RGBColor(0xAA, 0xBB, 0xCC)
TEXT_DIM    = RGBColor(0x55, 0x55, 0x55)

FONT = "Calibri"

# ── Slide dimensions — 20 × 11.25 in (WP3 template) ──────────────────────────
W        = Inches(20.0)
H        = Inches(11.25)
MARGIN   = Inches(0.75)
HEADER_H = Inches(1.55)
CY       = Inches(1.75)            # content top
CH       = H - CY - Inches(0.35)  # content height
CW       = W - 2 * MARGIN         # content width

# ── Paths ──────────────────────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).resolve().parent
ROOT       = SCRIPT_DIR.parents[1]   # workspace root
TEMPLATE   = ROOT / "05_documents_and_presentations" / "Role of Data and AI- WP3.pptx"
OUT_FILE   = ROOT / "05_documents_and_presentations" / "HospitalGuard_QuarterlyReview_May2026.pptx"


# ══════════════════════════════════════════════════════════════════════════════
# CORE HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def open_template() -> Presentation:
    """Open WP3 template and strip its existing slides, keeping theme/masters."""
    prs = Presentation(str(TEMPLATE))
    sldIdLst = prs.slides._sldIdLst
    for i in range(len(prs.slides) - 1, -1, -1):
        elem = sldIdLst[i]
        rId  = elem.get(qn("r:id"))
        prs.slides.part.drop_rel(rId)
        sldIdLst.remove(elem)
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def R(slide, x, y, w, h, fill: RGBColor,
      border: RGBColor = None, border_pt: float = 1.5) -> object:
    """Add a filled rectangle; optionally with a solid border."""
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(border_pt)
    else:
        sh.line.fill.background()
    return sh


def SR(slide, x, y, w, h, fill: RGBColor,
       border: RGBColor = None, border_pt: float = 1.5) -> object:
    """Drop-shadow rectangle (paint shadow first, then main rect on top)."""
    R(slide, x + Inches(0.09), y + Inches(0.09), w, h, SHADOW)
    return R(slide, x, y, w, h, fill, border, border_pt)


def T(slide, text: str, x, y, w, h,
      size: int = 18, bold: bool = False,
      color: RGBColor = NAVY,
      align=PP_ALIGN.LEFT, italic: bool = False) -> object:
    """Add a single-paragraph text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.name       = FONT
    run.font.size       = Pt(size)
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.color.rgb  = color
    return tb


def TM(slide, lines: list, x, y, w, h,
       size: int = 18, bold: bool = False,
       color: RGBColor = NAVY,
       align=PP_ALIGN.LEFT, line_spacing_pt: float = 6) -> object:
    """Multi-line text box (one paragraph per item in *lines*)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(line_spacing_pt)
        run = p.add_run()
        run.text           = line
        run.font.name      = FONT
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = color
    return tb


def HDR(slide, title: str, subtitle: str = ""):
    """Standard WP3 header: orange left bar + navy title + orange rule."""
    R(slide, 0, 0, Inches(0.22), HEADER_H, ORANGE)
    T(slide, title,
      MARGIN, Inches(0.18), W - MARGIN - Inches(0.3), Inches(0.9),
      size=44, bold=True, color=NAVY)
    R(slide, MARGIN, Inches(1.25), W - 2 * MARGIN, Inches(0.06), ORANGE)
    if subtitle:
        T(slide, subtitle,
          MARGIN, Inches(1.35), W - 2 * MARGIN, Inches(0.35),
          size=19, italic=True, color=TEXT_DIM)


def PAGE(slide, n: int, total: int = 7):
    """Slide number badge (bottom right)."""
    T(slide, f"{n:02d} / {total:02d}",
      W - Inches(1.8), H - Inches(0.45), Inches(1.55), Inches(0.38),
      size=13, color=MID_GREY, align=PP_ALIGN.RIGHT)


def STRIPE(slide):
    """Orange bottom rule + right dark sliver used on interior slides."""
    R(slide, 0, 0, Inches(0.22), H, ORANGE)
    R(slide, W - Inches(0.22), 0, Inches(0.22), H, DARK_NAVY2)
    R(slide, 0, H - Inches(0.08), W, Inches(0.08), ORANGE)


def PH(slide, x, y, w, h, media: str, description: str):
    """
    Dark placeholder block with dashed border.
    media:       "VIDEO" or "PHOTO"
    description: what the presenter should insert here
    """
    bg = R(slide, x, y, w, h, NEAR_BLACK, PH_BORDER, 2.0)
    # Set dashed border via XML
    try:
        spPr = bg._element.spPr
        ln_el = spPr.find(qn("a:ln"))
        if ln_el is None:
            ln_el = etree.SubElement(spPr, qn("a:ln"))
        if ln_el.find(qn("a:prstDash")) is None:
            pd = etree.SubElement(ln_el, qn("a:prstDash"))
            pd.set("val", "dashDot")
    except Exception:
        pass   # non-fatal: dashed style is cosmetic only

    icon = "▶  VIDEO" if media.upper() == "VIDEO" else "📷  PHOTO"
    T(slide, icon,
      x + Inches(0.25), y + Inches(0.25), w - Inches(0.5), Inches(0.45),
      size=16, bold=True, color=BLUE_STEEL, align=PP_ALIGN.LEFT)
    T(slide, description,
      x + Inches(0.25), y + h * 0.42, w - Inches(0.5), Inches(1.0),
      size=13, italic=True, color=TEXT_LIGHT, align=PP_ALIGN.LEFT)


def DOTS(slide, items: list, x, y, w, row_h,
         size: int = 19, dot_color: RGBColor = ORANGE,
         text_color: RGBColor = TEXT_DARK):
    """Bullet list with small coloured square indicators."""
    for i, text in enumerate(items):
        iy = y + i * row_h
        R(slide,
          x, iy + row_h * 0.35, Inches(0.18), Inches(0.18), dot_color)
        T(slide, text,
          x + Inches(0.32), iy, w - Inches(0.32), row_h,
          size=size, color=text_color)


def NOTES(slide, text: str):
    """Write full speaker notes to the notes pane."""
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.paragraphs[0].text = text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title & Executive Summary
# ══════════════════════════════════════════════════════════════════════════════
def slide1(prs):
    sl = blank(prs)

    # — Left navy panel --------------------------------------------------------
    R(sl, 0, 0, Inches(8.1), H, NAVY)
    R(sl, Inches(8.1), 0, Inches(0.25), H, ORANGE)

    # Slide badge
    R(sl, 0, H - Inches(0.55), Inches(1.2), Inches(0.45), ORANGE)
    T(sl, "01 / 07",
      Inches(0.04), H - Inches(0.55), Inches(1.15), Inches(0.45),
      size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Main title
    T(sl,
      "Robot Cognitive\nRecognition\nFramework",
      Inches(0.6), Inches(0.85), Inches(7.0), Inches(4.2),
      size=52, bold=True, color=WHITE)

    R(sl, Inches(0.6), Inches(5.45), Inches(5.2), Inches(0.1), ORANGE)

    T(sl, "Spatial 3D Mapping & Digital Twin",
      Inches(0.6), Inches(5.65), Inches(7.1), Inches(0.65),
      size=26, color=RGBColor(0xAA, 0xBB, 0xCC))

    T(sl, "Quarterly Review   ·   May 2026",
      Inches(0.6), Inches(6.45), Inches(7.1), Inches(0.5),
      size=22, color=TEXT_LIGHT)

    # — Key stats row (bottom of navy panel) -----------------------------------
    stats = [
        ("109",   "Object Classes"),
        ("3s",    "Occlusion Buffer"),
        ("X·Y·Z", "3D Coordinates"),
    ]
    sw = Inches(2.35)
    sy = H - Inches(2.85)
    for i, (val, lbl) in enumerate(stats):
        sx = Inches(0.28) + i * (sw + Inches(0.14))
        R(sl, sx, sy, sw, Inches(2.35), RGBColor(0x2A, 0x2C, 0x50))
        R(sl, sx, sy, sw, Inches(0.07), ORANGE)
        T(sl, val,
          sx, sy + Inches(0.15), sw, Inches(1.05),
          size=38, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        T(sl, lbl,
          sx, sy + Inches(1.2), sw, Inches(0.6),
          size=14, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # — Right panel placeholder ------------------------------------------------
    ph_x  = Inches(8.6)
    ph_w  = W - ph_x - Inches(0.35)
    ph_h  = H * 0.66
    ph_y  = H * 0.05
    PH(sl, ph_x, ph_y, ph_w, ph_h,
       "VIDEO",
       "INSERT: Robot live-demo video — the robot navigating a hospital\n"
       "or lab corridor while the system detects and labels objects\n"
       "in real time (annotated bounding boxes with class + ID).")

    # — Executive summary bullets below placeholder ----------------------------
    exec_y = ph_y + ph_h + Inches(0.3)
    exec_h = H - exec_y - Inches(0.25)
    exec_items = [
        "End-to-end pipeline now operational:  detect  →  track  →  map in 3D",
        "109 object categories  |  ByteTrack persistent IDs  |  SQLite Spatial Memory",
        "Off-board design keeps the robot agile while GPU stays on the base station",
    ]
    DOTS(sl, exec_items,
         ph_x + Inches(0.15), exec_y, ph_w - Inches(0.15),
         exec_h / 3, size=18,
         dot_color=ORANGE, text_color=NAVY)

    NOTES(sl,
          "Good morning / afternoon. This is our quarterly update on the Robot Cognitive Recognition Framework. "
          "What I want you to take away from this slide is simple: we have moved from separate experiments "
          "to a connected, working system. It can see 109 object categories, it tracks each one over time with "
          "a stable identity, and it knows where every tracked object is in 3D space. Those three capabilities "
          "together form the foundation of our hospital safety intelligence platform. "
          "On the right is a demo of the live system — feel free to pause on it as I speak.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Vision: Health & Safety in Hospitals
# ══════════════════════════════════════════════════════════════════════════════
def slide2(prs):
    sl = blank(prs)
    STRIPE(sl)
    HDR(sl,
        "The Vision — Health & Safety in Hospitals",
        "Why we are building this  ·  Three core safety pillars")
    PAGE(sl, 2)

    PILLAR = [
        (ORANGE,     "Clear Pathways",
         [
             "Detect obstructions in corridors in real time",
             "Flag blocked emergency routes immediately",
             "Support faster evacuation and emergency response",
         ],
         "PHOTO",
         "INSERT: Photo of a hospital corridor — ideally showing\n"
         "a wheelchair or equipment partially blocking the walkway\n"
         "(demonstrates the problem the system solves)."),

        (BLUE_STEEL, "Precision Alerting",
         [
             "Two-step verification before any alert is fired",
             "Dramatically fewer false alarms for busy staff",
             "Context-aware: sees the full scene, not just one object",
         ],
         "PHOTO",
         "INSERT: Photo of a nursing station monitor or alert\n"
         "screen — shows a staff member responding to an alert,\n"
         "conveying the human side of the alerting workflow."),

        (GREEN,      "Asset Tracking",
         [
             "Wheelchairs, IV stands, beds — always locatable",
             "Parked vs. moving state detected automatically",
             "Reduce time staff spend searching for equipment",
         ],
         "PHOTO",
         "INSERT: Photo of a hospital storeroom or ward bay\n"
         "with clearly labelled mobile medical equipment —\n"
         "wheelchairs, IV stands, patient monitors."),
    ]

    card_w = (CW - Inches(0.8)) / 3
    card_h = Inches(8.9)

    for i, (col, title, bullets, ph_type, ph_desc) in enumerate(PILLAR):
        cx = MARGIN + i * (card_w + Inches(0.4))
        cy = CY

        # Drop shadow + card
        SR(sl, cx, cy, card_w, card_h, WHITE, NAVY, 1.5)

        # Coloured top strip
        R(sl, cx, cy, card_w, Inches(0.85), col)

        # Pillar number badge
        R(sl, cx + Inches(0.2), cy + Inches(0.17), Inches(0.52), Inches(0.52),
          NEAR_BLACK)
        T(sl, str(i + 1),
          cx + Inches(0.2), cy + Inches(0.17), Inches(0.52), Inches(0.52),
          size=22, bold=True, color=col, align=PP_ALIGN.CENTER)

        T(sl, title,
          cx + Inches(0.85), cy + Inches(0.12),
          card_w - Inches(1.05), Inches(0.62),
          size=24, bold=True, color=WHITE)

        DOTS(sl, bullets,
             cx + Inches(0.3), cy + Inches(1.02),
             card_w - Inches(0.5), Inches(0.8),
             size=18, dot_color=col)

        # Photo placeholder
        PH(sl,
           cx + Inches(0.18), cy + Inches(3.65),
           card_w - Inches(0.36), Inches(4.9),
           ph_type, ph_desc)

    NOTES(sl,
          "The goal is practical and immediate: help hospitals run safer every day. "
          "Think of this as adding an always-alert safety observer who never gets tired and never gets distracted. "
          "Three things matter most to the hospital: first, are pathways clear for staff and emergency teams? "
          "Second, are alerts accurate? Nobody needs more noise in a busy ward. "
          "Third, is critical equipment where it should be? A lost wheelchair wastes five minutes of a nurse's time. "
          "These are not software features — they are direct contributions to patient and staff safety.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Cognitive Engine & Context Verification
# ══════════════════════════════════════════════════════════════════════════════
def slide3(prs):
    sl = blank(prs)
    STRIPE(sl)
    HDR(sl,
        "The Cognitive Engine — What We Detect & How We Stay Accurate",
        "YOLO V1 + V3 Ensemble  ·  Grounding DINO Fallback  ·  Two-Step Buddy System")
    PAGE(sl, 3)

    # — Left column: 3 stacked model cards with arrows -------------------------
    lx = MARGIN + Inches(0.1)
    lw = Inches(8.2)
    bh = Inches(2.45)
    gap = Inches(0.2)

    MODELS = [
        (NAVY,       "1",  "YOLO V1  —  106 Classes",
         "Covers all standard COCO + hospital objects. Runs on every single frame. "
         "Optimised for speed and broad coverage."),
        (BLUE_STEEL, "2",  "YOLO V3  —  109 Classes",
         "Adds bag, exit sign and spillage detection. Merged with V1 through "
         "ensemble logic to resolve overlapping detections cleanly."),
        (ORANGE,     "3",  "Grounding DINO  —  Fallback Layer",
         "Fires selectively for hard or rare classes (surgical scissors, IV stand, "
         "test tube). Context-gated: only activates when supporting objects are "
         "already confirmed in the scene."),
    ]

    for i, (col, num, title, body) in enumerate(MODELS):
        bx = lx
        by = CY + i * (bh + gap)

        SR(sl, bx, by, lw, bh, WHITE, col, 2.0)

        # Left accent strip with number
        R(sl, bx, by, Inches(0.65), bh, col)
        T(sl, num,
          bx + Inches(0.02), by + bh / 2 - Inches(0.45),
          Inches(0.62), Inches(0.9),
          size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        T(sl, title,
          bx + Inches(0.8), by + Inches(0.18),
          lw - Inches(1.0), Inches(0.6),
          size=22, bold=True, color=col)

        T(sl, body,
          bx + Inches(0.8), by + Inches(0.85),
          lw - Inches(1.0), Inches(1.4),
          size=17, color=TEXT_DARK)

        # Arrow between cards
        if i < len(MODELS) - 1:
            ay = by + bh + Inches(0.02)
            T(sl, "▼",
              bx + lw / 2 - Inches(0.35), ay,
              Inches(0.7), gap + Inches(0.02),
              size=14, bold=True, color=MID_GREY, align=PP_ALIGN.CENTER)

    # — Right column top: detection screenshot placeholder ----------------------
    rx  = lx + lw + Inches(0.55)
    rw  = W - rx - MARGIN - Inches(0.05)
    PH(sl, rx, CY, rw, Inches(4.7),
       "PHOTO",
       "INSERT: Screenshot from the live inference pipeline —\n"
       "annotated hospital scene showing bounding boxes with\n"
       "class labels and confidence % (e.g. wheelchair 94%, glove 83%).")

    # — Buddy system card (bottom right) ----------------------------------------
    by2 = CY + Inches(4.9)
    bcard_h = H - by2 - Inches(0.35)

    R(sl, rx, by2, rw, bcard_h, NAVY, ORANGE, 2.5)

    # Header
    R(sl, rx, by2, rw, Inches(0.75), RGBColor(0x2A, 0x2C, 0x50))
    T(sl, "THE BUDDY SYSTEM",
      rx + Inches(0.25), by2 + Inches(0.12),
      rw - Inches(0.5), Inches(0.5),
      size=20, bold=True, color=ORANGE)

    T(sl, "Two-Step Contextual Verification",
      rx + Inches(0.25), by2 + Inches(0.82),
      rw - Inches(0.5), Inches(0.4),
      size=15, bold=True, color=TEXT_LIGHT)

    RULES = [
        ("Test tube",        "only logged if  glove  or  medical tray  also seen"),
        ("Radiator",         "only logged if  wall  or  window  context present"),
        ("IV stand",         "only logged if  IV bag  or  patient bed  confirmed"),
        ("Surgical scissors","only logged if  glove  or  healthcare worker  near"),
    ]
    for j, (obj, rule) in enumerate(RULES):
        ry2 = by2 + Inches(1.3) + j * Inches(0.72)
        R(sl, rx + Inches(0.25), ry2 + Inches(0.22),
          Inches(0.14), Inches(0.14), ORANGE)
        T(sl, f"{obj}  →  {rule}",
          rx + Inches(0.5), ry2,
          rw - Inches(0.65), Inches(0.62),
          size=15, color=WHITE)

    NOTES(sl,
          "Our detection engine layers three systems. YOLO sweeps every frame, covering 109 classes. "
          "Where YOLO struggles — very small objects, unusual angles, rare equipment — Grounding DINO "
          "takes a selective second look. But the most important story here is accuracy without false alarms. "
          "We call it the Buddy System. Before the AI officially logs any uncertain detection, "
          "it must also find a believable companion in the same scene. "
          "If it thinks it sees a test tube, it must also see a medical tray or a glove. "
          "If it claims there is a radiator, there must be a wall or window context nearby. "
          "This two-step check eliminates hallucinations. "
          "In a hospital, false positives are not just annoying — they erode trust in the system. "
          "The Buddy System is how we ensure every logged detection is credible.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Advanced Tracking: The AI's Memory
# ══════════════════════════════════════════════════════════════════════════════
def slide4(prs):
    sl = blank(prs)
    STRIPE(sl)
    HDR(sl,
        "Advanced Tracking — The AI's Memory",
        "ByteTrack Name-Tag System  ·  3-Second Occlusion Tolerance  ·  Motion-Aware PPE Stabilizer")
    PAGE(sl, 4)

    CARDS = [
        (NAVY,       ORANGE,
         "NAME TAG SYSTEM",   "ByteTrack",
         [
             "Every person or object gets a persistent ID the moment it appears",
             "That ID follows the same entity across the entire camera frame",
             "One doctor walking right to left = Doctor #1, not two separate people",
         ],
         "Result:  Stable counts, zero double-counting"),

        (BLUE_STEEL, WHITE,
         "HIDE-AND-SEEK MEMORY",  "3-Second Occlusion Buffer",
         [
             "Nurse steps behind an IV pole — temporarily invisible",
             "System holds that identity file open for 3 seconds",
             "She steps out — same ID, no phantom reappear event",
         ],
         "Result:  No false disappear/reappear events"),

        (GREEN,      WHITE,
         "PPE STABILIZER",  "Motion-Aware Logic",
         [
             "Surgeon's arm sweeps across their face mask during a procedure",
             "System checks: is the associated body still moving and present?",
             "If yes — no 'Missing Mask!' alarm is triggered",
         ],
         "Result:  Fewer false alarms, higher compliance trust"),
    ]

    card_w = (CW - Inches(0.9)) / 3
    card_h = Inches(6.2)

    for i, (top_col, tag_col, icon_title, subtitle, bullets, result) in enumerate(CARDS):
        cx = MARGIN + i * (card_w + Inches(0.45))
        cy = CY

        SR(sl, cx, cy, card_w, card_h, WHITE, NAVY, 1.5)

        # Coloured header block
        R(sl, cx, cy, card_w, Inches(1.1), top_col)
        T(sl, icon_title,
          cx + Inches(0.18), cy + Inches(0.1),
          card_w - Inches(0.36), Inches(0.5),
          size=19, bold=True, color=WHITE)
        T(sl, subtitle,
          cx + Inches(0.18), cy + Inches(0.62),
          card_w - Inches(0.36), Inches(0.42),
          size=15, color=RGBColor(0xDD, 0xEE, 0xFF)
          if top_col == NAVY else RGBColor(0xEE, 0xFF, 0xEE))

        DOTS(sl, bullets,
             cx + Inches(0.28), cy + Inches(1.25),
             card_w - Inches(0.45), Inches(0.88),
             size=17, dot_color=top_col)

        # Result footer
        R(sl, cx, cy + card_h - Inches(0.65), card_w, Inches(0.65), top_col)
        T(sl, result,
          cx + Inches(0.18), cy + card_h - Inches(0.65),
          card_w - Inches(0.36), Inches(0.65),
          size=14, bold=True, color=WHITE)

    # Tracking video placeholder — below cards
    ph_y = CY + card_h + Inches(0.25)
    ph_h = H - ph_y - Inches(0.3)
    if ph_h > Inches(0.5):
        PH(sl, MARGIN, ph_y, CW, ph_h,
           "VIDEO",
           "INSERT: ByteTrack side-by-side demo video — left panel: raw camera feed, "
           "right panel: annotated feed showing persistent bounding boxes with "
           "ID numbers (e.g. #1 Nurse, #3 Wheelchair) following each entity smoothly "
           "across the scene even through brief occlusions.")

    NOTES(sl,
          "Let me bring tracking to life with three plain ideas. "
          "First, the Name Tag system. ByteTrack gives every detected person or object a persistent ID "
          "the moment they appear — like physically handing them a badge. "
          "One doctor walking left to right is always Doctor #1, not counted twice as two different people. "
          "Second, Hide-and-Seek Memory. If a nurse disappears behind an IV pole, "
          "we hold her file open for three seconds. When she steps back out, she is the same person. "
          "No phantom disappear-and-reappear event logged. "
          "Third, the PPE Stabilizer uses common sense: if a moving arm briefly covers a face mask, "
          "we check whether the body is still present before triggering any alarm. "
          "Together, these three layers mean the system produces calmer, more credible monitoring output.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Spatial 3D Mapping & Smart Filtering
# ══════════════════════════════════════════════════════════════════════════════
def slide5(prs):
    sl = blank(prs)
    STRIPE(sl)
    HDR(sl,
        "Spatial 3D Mapping — The Digital Twin",
        "Depth-fused 3D coordinates  ·  Semantic Spatial Memory  ·  Ghost Prevention Filtering")
    PAGE(sl, 5)

    lw = Inches(9.5)

    # — Top explainer card (full left-column width) ----------------------------
    R(sl, MARGIN, CY, lw, Inches(1.7), NAVY)
    T(sl,
      "Each tracked object  →  (X, Y, Z) + class name + track ID + timestamp",
      MARGIN + Inches(0.3), CY + Inches(0.2),
      lw - Inches(0.6), Inches(0.55),
      size=19, bold=True, color=ORANGE)
    T(sl,
      "The depth camera measures distance to every detected object. "
      "Combined with camera geometry (focal length, pixel position), "
      "we compute a real-world 3D coordinate and write it to the Spatial Memory database.",
      MARGIN + Inches(0.3), CY + Inches(0.8),
      lw - Inches(0.6), Inches(0.75),
      size=16, color=TEXT_LIGHT)

    # — Smart filtering table --------------------------------------------------
    T(sl, "SMART FILTERING  —  GHOST PREVENTION",
      MARGIN, CY + Inches(1.95), lw, Inches(0.5),
      size=19, bold=True, color=ORANGE)

    FILTER_ROWS = [
        (RGBColor(0xC0, 0x39, 0x2B),
         "DYNAMIC",   "People — staff, patients, visitors",
         "Lives in RAM only. Never written to the database permanently.",
         "Prevents the map filling with movement ghost-trails"),

        (BLUE_STEEL,
         "STATIC",    "Fixed assets — chairs, desks, monitors, doors",
         "Pinned to database immediately on first confirmed detection.",
         "Permanent 3D record for audit and spatial planning"),

        (ORANGE,
         "SEMI-STATIC", "Mobile assets — wheelchairs, beds, IV stands",
         "Uses the Parked Heuristic: velocity is monitored each frame.",
         "3D pin dropped only after 3 seconds of complete stillness"),
    ]

    row_h  = Inches(1.72)
    tbl_y  = CY + Inches(2.55)

    for i, (col, cat, desc, rule, result) in enumerate(FILTER_ROWS):
        ry = tbl_y + i * row_h

        SR(sl, MARGIN, ry, lw, row_h - Inches(0.1), WHITE, col, 1.8)
        R(sl, MARGIN, ry, Inches(0.6), row_h - Inches(0.1), col)

        # Category name in colour strip
        T(sl, cat,
          MARGIN + Inches(0.68), ry + Inches(0.1),
          Inches(2.2), Inches(0.5),
          size=17, bold=True, color=col)

        T(sl, desc,
          MARGIN + Inches(0.68), ry + Inches(0.6),
          lw - Inches(0.85), Inches(0.5),
          size=15, color=TEXT_DARK)

        T(sl, rule,
          MARGIN + Inches(0.68), ry + Inches(1.1),
          lw - Inches(1.1), Inches(0.42),
          size=14, color=TEXT_DARK)

        # Result strip
        R(sl, MARGIN + lw - Inches(5.1), ry + Inches(0.05),
          Inches(5.05), Inches(0.38), col)
        T(sl, "  " + result,
          MARGIN + lw - Inches(5.1), ry + Inches(0.05),
          Inches(5.05), Inches(0.38),
          size=13, bold=True, color=WHITE)

    # — Right column -----------------------------------------------------------
    rx = MARGIN + lw + Inches(0.5)
    rw = W - rx - MARGIN - Inches(0.05)

    # 3D map video placeholder
    PH(sl, rx, CY, rw, Inches(6.5),
       "VIDEO",
       "INSERT: 3D Spatial Map visualization — top-down or perspective\n"
       "ward floor plan showing coloured 3D pins for wheelchairs,\n"
       "fire extinguishers, IV stands updating in real-time as the\n"
       "robot moves (ideally from hospital_twin.db export or RViz).")

    # — Sample DB record card --------------------------------------------------
    db_y = CY + Inches(6.7)
    db_h = H - db_y - Inches(0.3)
    R(sl, rx, db_y, rw, db_h, NAVY)
    T(sl, "Sample Spatial Memory Record",
      rx + Inches(0.2), db_y + Inches(0.12),
      rw - Inches(0.4), Inches(0.4),
      size=14, bold=True, color=ORANGE)
    TM(sl,
       [
           "class        |  ID  |   X      Y       Z    |  last_seen",
           "wheelchair   |  12  | +1.42  −0.31  +2.85  | 10:14:06",
           "fire_ext     |   3  | −2.10  +0.85  +4.12  | 10:14:04",
           "patient      |  44  | +0.22  −0.04  +1.97  | 10:14:07",
       ],
       rx + Inches(0.2), db_y + Inches(0.6),
       rw - Inches(0.35), db_h - Inches(0.7),
       size=12, color=TEXT_LIGHT)

    NOTES(sl,
          "Detection tells us what is present. Tracking gives it a persistent identity. "
          "Spatial mapping tells us exactly where it is in the room. "
          "We use depth camera data to back-calculate X, Y, Z coordinates for every tracked object "
          "and write them to the hospital twin database with a timestamp. "
          "But here is the engineering challenge: if we saved everything, the map would fill with "
          "ghost trails of every person who walked through. We call this Ghost Prevention. "
          "People stay in live memory only — never permanently saved. "
          "Fixed assets like chairs are written to the database immediately. "
          "Semi-static assets like wheelchairs use the Parked Heuristic: "
          "we monitor velocity, and only when the object has sat completely still for three seconds "
          "do we drop a permanent 3D pin. "
          "The result is a clean, actionable spatial map rather than a record full of clutter.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Off-Board Processing Architecture
# ══════════════════════════════════════════════════════════════════════════════
def slide6(prs):
    sl = blank(prs)
    STRIPE(sl)
    HDR(sl,
        "Off-Board Processing Architecture",
        "Robot = sensor node  ·  Base station = compute engine  ·  Clean data separation")
    PAGE(sl, 6)

    # — Architecture flow diagram: 5 nodes ------------------------------------
    NODES = [
        (NAVY,                  "ROBOT\nCHASSIS",
         "Orbbec RGB-D Camera\nLiDAR Sensor\nWi-Fi Module",
         "Sensor Node Only"),

        (RGBColor(0x25, 0x27, 0x55), "Wi-Fi\nBRIDGE",
         "Compressed stream\nUDP / TCP tunnel\nPacket-loss recovery",
         "Connectivity Layer"),

        (BLUE_STEEL,            "BASE\nSTATION",
         "YOLO + DINO Inference\nByteTrack Processing\n3D Coord. Calculation",
         "Compute Engine"),

        (GREEN,                 "SPATIAL\nDATABASE",
         "SQLite Memory Store\nTimestamped 3D Records\nAudit Trail",
         "Persistent Store"),

        (ORANGE,                "OUTPUTS",
         "Annotated Video Feed\nReal-time Safety Alerts\nCompliance Dashboard",
         "Operational Value"),
    ]

    n      = len(NODES)
    gap    = Inches(0.32)
    node_w = (CW - gap * (n - 1)) / n
    flow_y = CY + Inches(0.1)
    flow_h = Inches(3.65)

    for i, (col, title, body, tag) in enumerate(NODES):
        nx = MARGIN + i * (node_w + gap)

        SR(sl, nx, flow_y, node_w, flow_h, NAVY, col, 2.5)

        # Coloured top cap
        R(sl, nx, flow_y, node_w, Inches(0.95), col)
        T(sl, title,
          nx + Inches(0.1), flow_y + Inches(0.08),
          node_w - Inches(0.2), Inches(0.82),
          size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        T(sl, body,
          nx + Inches(0.18), flow_y + Inches(1.05),
          node_w - Inches(0.36), Inches(2.1),
          size=14, color=TEXT_LIGHT)

        # Tag footer
        R(sl, nx, flow_y + flow_h - Inches(0.58), node_w, Inches(0.58), col)
        T(sl, tag,
          nx, flow_y + flow_h - Inches(0.58), node_w, Inches(0.58),
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Arrow connector
        if i < n - 1:
            T(sl, "→",
              nx + node_w + Inches(0.04), flow_y + flow_h / 2 - Inches(0.3),
              gap - Inches(0.04), Inches(0.55),
              size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # — Two Why-cards below the flow -------------------------------------------
    why_y = flow_y + flow_h + Inches(0.5)
    why_h = H - why_y - Inches(0.35)
    why_w = (CW - Inches(0.5)) / 2

    WHY = [
        (NAVY,       "Why This Architecture?",
         [
             "Robot stays light and agile — no heavy GPU on board",
             "AI models upgraded centrally, no hardware changes on robot",
             "Higher frame rates: inference doesn't compete with mobility",
             "Scales to multiple rooms by adding cameras to one station",
         ]),
        (BLUE_STEEL, "Active Engineering Challenges",
         [
             "Wi-Fi packet-loss recovery algorithm under development",
             "Compression pipeline tuned for <100 ms round-trip target",
             "Multi-device concurrent stream stress-testing in progress",
             "Failsafe buffering to maintain tracking through dropouts",
         ]),
    ]

    for i, (col, title, pts) in enumerate(WHY):
        wx = MARGIN + i * (why_w + Inches(0.5))
        SR(sl, wx, why_y, why_w, why_h, NEAR_BLACK, col, 2.0)

        R(sl, wx, why_y, why_w, Inches(0.65), col)
        T(sl, title,
          wx + Inches(0.3), why_y + Inches(0.1),
          why_w - Inches(0.6), Inches(0.5),
          size=19, bold=True, color=WHITE)

        DOTS(sl, pts,
             wx + Inches(0.3), why_y + Inches(0.8),
             why_w - Inches(0.5), (why_h - Inches(0.9)) / 4,
             size=16, dot_color=ORANGE,
             text_color=RGBColor(0xCC, 0xDD, 0xEE))

    NOTES(sl,
          "We made a deliberate design decision: keep the robot simple and put the heavy computation elsewhere. "
          "The robot focuses on collecting sensor data and streaming it over Wi-Fi. "
          "The base station runs the YOLO models, manages ByteTrack, calculates 3D coordinates, "
          "and writes to the database. "
          "Think of it like a camera crew and a post-production studio: the camera captures, the studio processes. "
          "This separation means we can upgrade the AI models centrally without touching robot hardware, "
          "maintain higher frame rates because inference doesn't compete with movement processing, "
          "and scale to multiple rooms just by pointing additional cameras at the same base station. "
          "The engineering focus right now is hardening the Wi-Fi bridge: "
          "handling packet loss, reducing latency, and validating stable stream quality "
          "under realistic hospital network conditions.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Next Phase & Future Deployment
# ══════════════════════════════════════════════════════════════════════════════
def slide7(prs):
    sl = blank(prs)
    STRIPE(sl)
    HDR(sl,
        "Next Phase — Towards Live Deployment",
        "Network hardening  ·  Untethered robot trial  ·  Hospital pilot preparation")
    PAGE(sl, 7)

    MILESTONES = [
        (ORANGE,     "PHASE 1",   "Network Hardening",
         [
             "Stress-test wireless bridge under realistic hospital network load",
             "Implement packet-loss recovery and jitter-buffering algorithms",
             "Validate continuous stream quality for stable YOLO inference",
             "Benchmark: sustained 25+ fps with <5 % packet loss threshold",
         ],
         "STATUS: IN PROGRESS"),

        (BLUE_STEEL, "PHASE 2",   "Untethered Robot Trial",
         [
             "First live end-to-end run on the physical robot chassis",
             "Sensor → Wi-Fi → inference → 3D map — fully untethered",
             "Capture performance baseline and safety check results",
             "Iterate on sensor calibration and tracking stability",
         ],
         "TARGET: Next Quarter"),

        (GREEN,      "PHASE 3",   "Controlled Hospital Pilot",
         [
             "Deploy in a single controlled ward environment",
             "Baseline KPIs: alert accuracy, asset tracking rate, pathway events",
             "Iterative feedback loop with clinical operations staff",
             "Gate review: readiness criteria before wider roll-out",
         ],
         "TARGET: Q3 2026"),
    ]

    mw     = (CW - Inches(1.0)) / 3
    mh     = Inches(6.2)
    mile_y = CY

    for i, (col, phase, title, pts, status) in enumerate(MILESTONES):
        mx = MARGIN + i * (mw + Inches(0.5))

        SR(sl, mx, mile_y, mw, mh, WHITE, col, 2.5)

        # Header
        R(sl, mx, mile_y, mw, Inches(1.2), col)
        T(sl, phase,
          mx + Inches(0.18), mile_y + Inches(0.06),
          mw - Inches(0.36), Inches(0.42),
          size=13, bold=True, color=WHITE)
        T(sl, title,
          mx + Inches(0.18), mile_y + Inches(0.5),
          mw - Inches(0.36), Inches(0.62),
          size=21, bold=True, color=WHITE)

        DOTS(sl, pts,
             mx + Inches(0.25), mile_y + Inches(1.35),
             mw - Inches(0.4), Inches(1.05),
             size=17, dot_color=col)

        # Status badge
        R(sl, mx + Inches(0.22), mile_y + mh - Inches(0.75),
          mw - Inches(0.44), Inches(0.62), col)
        T(sl, status,
          mx + Inches(0.22), mile_y + mh - Inches(0.75),
          mw - Inches(0.44), Inches(0.62),
          size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Connector arrow between milestones
        if i < len(MILESTONES) - 1:
            T(sl, "→",
              mx + mw + Inches(0.05), mile_y + mh / 2 - Inches(0.3),
              Inches(0.45), Inches(0.55),
              size=24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # — Bottom row: photo placeholder + closing statement ----------------------
    ph_y = mile_y + mh + Inches(0.3)
    ph_h = H - ph_y - Inches(0.3)
    ph_w = Inches(8.5)

    if ph_h > Inches(0.45):
        PH(sl, MARGIN, ph_y, ph_w, ph_h,
           "PHOTO",
           "INSERT: Photo of the physical robot chassis in the lab —\n"
           "ideally with the Orbbec RGB-D camera mounted and\n"
           "the team setting it up for the first untethered test.")

        # Closing statement card
        close_x = MARGIN + ph_w + Inches(0.4)
        close_w = CW - ph_w - Inches(0.4)
        R(sl, close_x, ph_y, close_w, ph_h, NAVY)
        R(sl, close_x, ph_y, close_w, Inches(0.07), ORANGE)

        T(sl,
          "The system works in the lab.\nThe next step is proving it\nworks in the world.",
          close_x + Inches(0.3), ph_y + Inches(0.2),
          close_w - Inches(0.55), ph_h - Inches(0.8),
          size=23, bold=True, color=WHITE)

        T(sl, "Questions welcome.",
          close_x + Inches(0.3), ph_y + ph_h - Inches(0.65),
          close_w - Inches(0.55), Inches(0.5),
          size=18, italic=True, color=ORANGE)

    NOTES(sl,
          "We have a working system. The next phase is about proving it works outside the lab. "
          "Phase one is network hardening: we are stress-testing the wireless bridge to ensure "
          "the data stream from the robot remains reliable under hospital network conditions — "
          "handling packet loss and jitter without breaking tracking or 3D mapping quality. "
          "Phase two is the first fully untethered live test on the physical robot chassis: "
          "sensor data, inference, tracking, and spatial mapping all running in real time as the robot moves. "
          "Phase three, our target for Q3 2026, is a controlled pilot in a real ward environment "
          "where we baseline our key performance indicators and start iterating with clinical staff. "
          "The system works in the lab. The next step is proving it works in the world. "
          "I am happy to take any questions now.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    print(f"Opening template: {TEMPLATE}")
    prs = open_template()

    print("Building Slide 1 — Title & Executive Summary …")
    slide1(prs)
    print("Building Slide 2 — Vision: Health & Safety …")
    slide2(prs)
    print("Building Slide 3 — Cognitive Engine …")
    slide3(prs)
    print("Building Slide 4 — Advanced Tracking …")
    slide4(prs)
    print("Building Slide 5 — Spatial 3D Mapping …")
    slide5(prs)
    print("Building Slide 6 — Off-Board Architecture …")
    slide6(prs)
    print("Building Slide 7 — Next Phase …")
    slide7(prs)

    OUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_FILE))
    print(f"\nSaved: {OUT_FILE}")


if __name__ == "__main__":
    main()
