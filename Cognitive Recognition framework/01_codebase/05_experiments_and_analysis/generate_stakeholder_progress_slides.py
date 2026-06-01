"""
Generate a non-technical stakeholder slide deck for HospitalGuard progress.

Usage:
  python generate_stakeholder_progress_slides.py

Output:
  05_documents_and_presentations/HospitalGuard_Stakeholder_Progress_May2026.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# Presentation sizing: 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Color theme
NAVY = RGBColor(0x11, 0x2A, 0x46)
TEAL = RGBColor(0x00, 0x8A, 0x9C)
SKY = RGBColor(0xD8, 0xF2, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1F, 0x2D, 0x3A)
MUTED = RGBColor(0x56, 0x66, 0x76)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
ORANGE = RGBColor(0xF2, 0x8E, 0x2B)

FONT = "Calibri"


def add_bg(slide, color=WHITE):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_header(slide, title: str, subtitle: str = ""):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.35), Inches(0.15), Inches(8.9), Inches(0.55))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sx = slide.shapes.add_textbox(Inches(9.2), Inches(0.2), Inches(3.8), Inches(0.45))
        sf = sx.text_frame
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = FONT
        sp.font.size = Pt(12)
        sp.font.color.rgb = SKY
        sp.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, left, top, width, height, items, size=20):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)


def add_card(slide, x, y, w, h, title, lines, accent=TEAL):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent
    card.line.width = Pt(1.5)

    title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.35))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent

    add_bullets(
        slide,
        x + Inches(0.2),
        y + Inches(0.55),
        w - Inches(0.35),
        h - Inches(0.65),
        lines,
        size=14,
    )


def add_flow_box(slide, x, y, w, h, text, fill=SKY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = TEAL
    box.line.width = Pt(1.25)

    tx = slide.shapes.add_textbox(x + Inches(0.08), y + Inches(0.08), w - Inches(0.16), h - Inches(0.16))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1
    return box


def connect_right_arrow(slide, from_shape, to_shape):
    x1 = from_shape.left + from_shape.width
    y1 = from_shape.top + int(from_shape.height / 2)
    x2 = to_shape.left
    y2 = to_shape.top + int(to_shape.height / 2)

    arrow = slide.shapes.add_connector(1, x1, y1, x2, y2)
    arrow.line.color.rgb = TEAL
    arrow.line.width = Pt(1.5)
    arrow.line.end_arrowhead = True


def build_deck() -> Path:
    root = Path(__file__).resolve().parents[2]
    out_path = root / "05_documents_and_presentations" / "HospitalGuard_Stakeholder_Progress_May2026.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # Slide 1: Title
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_header(s, "HospitalGuard Progress Update", "May 2026")

    title = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12.0), Inches(1.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "AI for Hospital Health and Safety"
    p.font.name = FONT
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = NAVY

    sub = s.shapes.add_textbox(Inches(0.75), Inches(3.0), Inches(11.8), Inches(1.0))
    sf = sub.text_frame
    sp = sf.paragraphs[0]
    sp.text = "From object detection to tracked, spatially anchored safety intelligence"
    sp.font.name = FONT
    sp.font.size = Pt(23)
    sp.font.color.rgb = MUTED

    metrics = [
        "109 object categories supported",
        "Real-time tracking with persistent IDs",
        "3D spatial memory stored in SQLite",
        "Designed for off-board processing pipelines",
    ]
    add_bullets(s, Inches(0.9), Inches(4.2), Inches(12.0), Inches(2.5), metrics, size=20)

    # Slide 2: Progress so far
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Progress Achieved")
    progress = [
        "Established a clean, production-oriented framework with separated training, inference, data, outputs, and docs.",
        "Built and validated HospitalGuard model variants (V1 + V3) for hospital-specific detection coverage.",
        "Integrated temporal tracking using ByteTrack for stable IDs across frames.",
        "Delivered RGB-D pipeline with 3D coordinate assignment per tracked object.",
        "Implemented long-term spatial memory logging to a hospital twin database.",
        "Prepared replay and export workflows for demos, validation, and stakeholder reporting.",
    ]
    add_bullets(s, Inches(0.8), Inches(1.2), Inches(12.0), Inches(5.9), progress, size=19)

    # Slide 3: Models used
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Models Used in Latest Architecture")

    add_card(
        s,
        Inches(0.6),
        Inches(1.3),
        Inches(4.1),
        Inches(2.4),
        "YOLO V1 (106 classes)",
        [
            "Base detector for broad object coverage",
            "Strong on general and hospital classes",
            "Runs every frame",
        ],
        accent=NAVY,
    )

    add_card(
        s,
        Inches(4.95),
        Inches(1.3),
        Inches(4.1),
        Inches(2.4),
        "YOLO V3 (109 classes)",
        [
            "Adds and reinforces priority classes",
            "Contributes overlap and new-class signals",
            "Merged through ensemble routing",
        ],
        accent=TEAL,
    )

    add_card(
        s,
        Inches(9.3),
        Inches(1.3),
        Inches(3.45),
        Inches(2.4),
        "Grounding DINO",
        [
            "Fallback for weak or hard classes",
            "Context gated to lower false positives",
            "Used selectively, not always-on",
        ],
        accent=ORANGE,
    )

    add_card(
        s,
        Inches(0.8),
        Inches(4.0),
        Inches(11.9),
        Inches(2.8),
        "What this means for stakeholders",
        [
            "The system combines speed (YOLO), resilience (ensemble), and recovery on difficult cases (DINO fallback).",
            "This design improves reliability in realistic hospital scenes where lighting, occlusion, and object size vary.",
        ],
        accent=GREEN,
    )

    # Slide 4: Architecture flow
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Latest Architecture: End-to-End Flow")

    b1 = add_flow_box(s, Inches(0.35), Inches(2.0), Inches(2.2), Inches(1.4), "RGB / RGB-D\nInput Stream")
    b2 = add_flow_box(s, Inches(2.75), Inches(2.0), Inches(2.1), Inches(1.4), "YOLO V1 + V3\nEnsemble")
    b3 = add_flow_box(s, Inches(5.05), Inches(2.0), Inches(2.1), Inches(1.4), "Grounding DINO\nFallback")
    b4 = add_flow_box(s, Inches(7.35), Inches(2.0), Inches(2.1), Inches(1.4), "ByteTrack\nTemporal IDs")
    b5 = add_flow_box(s, Inches(9.65), Inches(2.0), Inches(3.35), Inches(1.4), "Spatial Memory\n(X,Y,Z + class + track_id)")

    connect_right_arrow(s, b1, b2)
    connect_right_arrow(s, b2, b3)
    connect_right_arrow(s, b3, b4)
    connect_right_arrow(s, b4, b5)

    add_bullets(
        s,
        Inches(0.8),
        Inches(4.0),
        Inches(12.0),
        Inches(2.7),
        [
            "Each detected object gets a consistent ID over time and a spatial coordinate in 3D space.",
            "This turns raw video into an operational memory that can support alerts, audits, and workflow analytics.",
        ],
        size=18,
    )

    # Slide 5: ByteTrack for non-technical audience
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "ByteTrack: Why It Matters")

    add_card(
        s,
        Inches(0.7),
        Inches(1.4),
        Inches(5.95),
        Inches(5.6),
        "Without tracking",
        [
            "The same person or object is counted repeatedly.",
            "Brief occlusions can appear as disappear-reappear events.",
            "Alerts become noisy and less trusted.",
        ],
        accent=ORANGE,
    )

    add_card(
        s,
        Inches(6.8),
        Inches(1.4),
        Inches(5.85),
        Inches(5.6),
        "With ByteTrack",
        [
            "Objects keep a stable identity while moving through the scene.",
            "Short occlusions are handled, reducing duplicate events.",
            "Enables time-aware logic: dwell, path, and safety behavior patterns.",
            "Creates a bridge from detection to memory and actionable analytics.",
        ],
        accent=GREEN,
    )

    # Slide 6: 3D map and coordinate memory
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Spatial 3D Map and Object Coordinates")

    add_bullets(
        s,
        Inches(0.8),
        Inches(1.25),
        Inches(12.0),
        Inches(1.45),
        [
            "RGB-D depth is fused with detections to assign world coordinates (X, Y, Z) per tracked object.",
            "Coordinates are written to the spatial memory table with timestamp, class name, and track ID.",
        ],
        size=18,
    )

    table_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.9), Inches(11.8), Inches(3.4))
    table_box.fill.solid()
    table_box.fill.fore_color.rgb = SKY
    table_box.line.color.rgb = TEAL

    tbl_title = s.shapes.add_textbox(Inches(1.1), Inches(3.1), Inches(11.3), Inches(0.35))
    tt = tbl_title.text_frame.paragraphs[0]
    tt.text = "Spatial Memory Record (conceptual)"
    tt.font.name = FONT
    tt.font.bold = True
    tt.font.size = Pt(16)
    tt.font.color.rgb = NAVY

    concept_rows = [
        "timestamp | class_name | track_id | X | Y | Z | last_seen",
        "2026-05-29T10:14:03 | wheelchair | 12 | 1.42 | -0.31 | 2.85 | 2026-05-29T10:14:06",
        "2026-05-29T10:14:04 | fire_extinguisher | 3 | -2.10 | 0.85 | 4.12 | 2026-05-29T10:14:04",
        "2026-05-29T10:14:05 | patient | 44 | 0.22 | -0.04 | 1.97 | 2026-05-29T10:14:07",
    ]
    add_bullets(s, Inches(1.15), Inches(3.55), Inches(11.2), Inches(2.55), concept_rows, size=13)

    # Slide 7: Memory behavior by object type
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "How Memory Differs by Object Type")

    add_card(
        s,
        Inches(0.6),
        Inches(1.35),
        Inches(4.1),
        Inches(5.6),
        "Static infrastructure",
        [
            "Examples: exit sign, fire extinguisher, doors",
            "Expected to remain in fixed areas",
            "Useful for compliance and missing-object checks",
        ],
        accent=NAVY,
    )

    add_card(
        s,
        Inches(4.95),
        Inches(1.35),
        Inches(4.1),
        Inches(5.6),
        "Mobile equipment",
        [
            "Examples: wheelchairs, IV stands, infusion pumps",
            "Track movement and utilization over time",
            "Supports asset availability and bottleneck analysis",
        ],
        accent=TEAL,
    )

    add_card(
        s,
        Inches(9.3),
        Inches(1.35),
        Inches(3.45),
        Inches(5.6),
        "People and PPE",
        [
            "Examples: staff, patients, mask, glove, hair net",
            "Higher variability and occlusion",
            "Combines short-term tracking with room-state memory",
        ],
        accent=ORANGE,
    )

    # Slide 8: Off-board processing
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Off-Board Processing Strategy")

    left = [
        "Edge camera captures RGB or RGB-D feeds.",
        "Detection and tracking can run on a dedicated GPU workstation/server.",
        "Only events, trajectories, and metadata are persisted centrally.",
        "Spatial memory database supports audit trails and post-event review.",
    ]
    right = [
        "Lower device complexity at the point of care",
        "Easier upgrades of models without replacing cameras",
        "Centralized governance and quality control",
        "Better scalability across wards and facilities",
    ]

    add_card(s, Inches(0.7), Inches(1.45), Inches(5.9), Inches(5.45), "How it works", left, accent=TEAL)
    add_card(s, Inches(6.75), Inches(1.45), Inches(5.85), Inches(5.45), "Business value", right, accent=GREEN)

    # Slide 9: Health and safety outcomes
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Hospital Health and Safety Impact")

    outcomes = [
        "PPE awareness: supports visibility of mask, glove, and hair-net related behaviors in clinical zones.",
        "Hazard detection: spots spillage and safety equipment such as fire extinguishers and exit signs.",
        "Equipment readiness: tracks key assets like wheelchairs and IV equipment in time and space.",
        "Operational visibility: creates explainable event history tied to object IDs and 3D coordinates.",
        "Decision support: enables targeted interventions with fewer blind spots than manual monitoring.",
    ]
    add_bullets(s, Inches(0.8), Inches(1.3), Inches(12.0), Inches(4.7), outcomes, size=19)

    close = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.05), Inches(12.0), Inches(0.95))
    close.fill.solid()
    close.fill.fore_color.rgb = NAVY
    close.line.fill.background()
    cbox = s.shapes.add_textbox(Inches(1.0), Inches(6.25), Inches(11.6), Inches(0.55))
    cp = cbox.text_frame.paragraphs[0]
    cp.text = "Next phase: pilot deployment, KPI baselining, and workflow-specific alert tuning"
    cp.font.name = FONT
    cp.font.size = Pt(18)
    cp.font.bold = True
    cp.font.color.rgb = WHITE
    cp.alignment = PP_ALIGN.CENTER

    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    generated = build_deck()
    print(f"Saved stakeholder deck to: {generated}")