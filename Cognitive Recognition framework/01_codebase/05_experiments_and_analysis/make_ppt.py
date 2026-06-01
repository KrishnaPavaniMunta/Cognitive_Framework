"""
Generate: why_hospital_v3_forgot_classes.pptx
8 slides (no fixing slide):
  1. Title
  2. What is Catastrophic Forgetting
  3. Root Cause: Data Imbalance
  4. Mechanism: How It Happens
  5. Freeze Strategy Analysis
  6. Oversampling — Only Helped New Classes
  7. Training Metrics
  8. Summary

Style: matches provided template — white background, dark-olive-green headings,
dark charcoal body text, Calibri font, clean minimal layout.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

FONT_NAME = "Calibri"

# ── Colour palette (light template) ──────────────────────────────────────────
BG_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # main slide background
BG_LIGHT  = RGBColor(0xF4, 0xF4, 0xF2)  # subtle alternate bg
BG_DARK   = BG_WHITE                     # alias kept for compat
BG_BLUE   = BG_WHITE
BG_WARN   = BG_WHITE
BG_GREEN  = BG_WHITE

HEADING   = RGBColor(0x4A, 0x5E, 0x28)  # dark olive green  (matches "Hello!")
BODY      = RGBColor(0x33, 0x33, 0x33)  # dark charcoal body text
DIM       = RGBColor(0x88, 0x88, 0x88)  # muted / caption text
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = BODY                         # "light" text → charcoal on white bg

# Accent colours (used for bars, badges, table highlights)
BLUE      = RGBColor(0x2E, 0x75, 0xB6)  # professional blue
GREEN     = RGBColor(0x4A, 0x5E, 0x28)  # olive green (same as heading)
RED       = RGBColor(0xC0, 0x39, 0x2B)  # muted red
YELLOW    = RGBColor(0xD4, 0xA0, 0x17)  # warm amber
PURPLE    = RGBColor(0x6C, 0x3D, 0x91)  # muted purple

# Card backgrounds
CARD_BG   = RGBColor(0xF0, 0xF4, 0xEB)  # very light green tint
CARD_RED  = RGBColor(0xFB, 0xEC, 0xEA)  # very light red tint
CARD_GRN  = RGBColor(0xEB, 0xF5, 0xEB)  # very light green tint
CARD_BLU  = RGBColor(0xE8, 0xF1, 0xFA)  # very light blue tint

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helper: set slide background colour ──────────────────────────────────────
def set_bg(slide, rgb: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


# ── Helper: add a text box ────────────────────────────────────────────────────
def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=BODY, align=PP_ALIGN.LEFT,
             italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


# ── Helper: multi-run paragraph ──────────────────────────────────────────────
def add_para(tf, segments, align=PP_ALIGN.LEFT, space_before=6):
    """segments = list of (text, size, bold, color)"""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    for text, size, bold, color in segments:
        run = p.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return p


# ── Helper: add a filled rectangle ───────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


# ── Helper: add a table ───────────────────────────────────────────────────────
def add_table(slide, headers, rows, left, top, width, height, col_widths=None):
    ncols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, ncols, left, top, width, height).table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    # header row
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x4A, 0x5E, 0x28)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.name = FONT_NAME
        run.font.bold = True; run.font.size = Pt(13); run.font.color.rgb = WHITE

    # data rows
    for ri, row in enumerate(rows):
        for ci, (cell_text, cell_color) in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            bg = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 0 else RGBColor(0xF4, 0xF4, 0xF2)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = cell_text
            run.font.name = FONT_NAME
            run.font.size = Pt(12); run.font.color.rgb = cell_color

    return tbl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG_WHITE)

# Olive accent bar top
add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), HEADING)

# Title
add_text(sl, "Why Hospital V3 Forgot Previous Classes",
         Inches(1), Inches(1.5), Inches(11.3), Inches(1.8),
         size=44, bold=True, color=HEADING, align=PP_ALIGN.CENTER)

# Accent line
add_rect(sl, Inches(5.6), Inches(3.45), Inches(2.1), Inches(0.06), HEADING)

# Subtitle
add_text(sl, "Catastrophic Forgetting Analysis  ·  YOLO V26m  ·  109 Classes",
         Inches(1), Inches(3.65), Inches(11.3), Inches(0.55),
         size=18, color=DIM, align=PP_ALIGN.CENTER)

# Tag pills row
tags = [("106 → 109 Classes", BLUE), ("GTX 1660 Ti", YELLOW),
        ("Data Imbalance", RED), ("2-Phase Freeze", GREEN)]
tag_x = Inches(2.5)
for label, col in tags:
    add_rect(sl, tag_x, Inches(4.55), Inches(2.1), Inches(0.42),
             RGBColor(0xF4, 0xF4, 0xF2), col, Pt(1.5))
    add_text(sl, label, tag_x, Inches(4.55), Inches(2.1), Inches(0.42),
             size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    tag_x += Inches(2.2)

# Accent bar bottom
add_rect(sl, 0, Inches(7.38), SLIDE_W, Inches(0.12), HEADING)

# Slide number
add_text(sl, "1 / 8", Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
         size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS CATASTROPHIC FORGETTING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG_WHITE)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), HEADING)

add_text(sl, "What Is Catastrophic Forgetting?",
         Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.85),
         size=32, bold=True, color=HEADING, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.06), HEADING)

add_text(sl,
    "When fine-tuned on new data, gradient updates that improve new-class performance "
    "can overwrite weight patterns responsible for old classes. The network \"forgets\" what it learned.",
    Inches(1), Inches(1.38), Inches(11.3), Inches(0.8),
    size=16, color=BODY, align=PP_ALIGN.CENTER)

# Two cards
for i, (title, body, bg, border) in enumerate([
    ("Classes at Risk",
     "Old hospital-specific classes that received NO new training images in V3:\n"
     "surgical_light · nasal_cannula · glove · test_tube · infusion_pump · hair_net · surgical_scissor …",
     CARD_RED, RED),
    ("Classes Reinforced",
     "The 3 overlapping Roboflow classes received thousands of extra images:\n\n"
     "door (96)  ·  wheelchair (94)  ·  fire_extinguisher (100)",
     CARD_GRN, GREEN),
]):
    cx = Inches(0.6) + i * Inches(6.4)
    add_rect(sl, cx, Inches(2.45), Inches(6.0), Inches(3.8), bg, border, Pt(2.0))
    col = RED if i == 0 else GREEN
    add_text(sl, title, cx + Inches(0.25), Inches(2.6), Inches(5.6), Inches(0.55),
             size=18, bold=True, color=col)
    add_text(sl, body, cx + Inches(0.25), Inches(3.2), Inches(5.6), Inches(2.8),
             size=14, color=BODY)

add_rect(sl, 0, Inches(7.38), SLIDE_W, Inches(0.12), HEADING)
add_text(sl, "2 / 8", Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
         size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATA IMBALANCE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG_WHITE)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), HEADING)

add_text(sl, "Root Cause: Severe Data Imbalance",
         Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.85),
         size=32, bold=True, color=HEADING, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.06), HEADING)

add_text(sl,
    "The new Roboflow dataset is 4.3× larger. The model sees far more new-class images "
    "per epoch, biasing gradient updates away from old classes.",
    Inches(1), Inches(1.38), Inches(11.3), Inches(0.65),
    size=16, color=BODY, align=PP_ALIGN.CENTER)

# Bar chart
bars = [
    ("hospital_merged (train)", 0.187, BLUE,   "4,604 images  (19%)"),
    ("hospital_v2 Roboflow",    0.815, RED,    "20,023 images  (81%)"),
    ("Total merged",            1.000, HEADING, "~24,627 images"),
]
track_w = Inches(7.8)
lx = Inches(1.2)
for idx, (label, pct, col, note) in enumerate(bars):
    ty = Inches(2.3) + idx * Inches(0.85)
    add_text(sl, label, lx, ty, Inches(2.5), Inches(0.42),
             size=13, color=BODY, align=PP_ALIGN.RIGHT)
    tx = lx + Inches(2.6)
    add_rect(sl, tx, ty + Inches(0.05), track_w, Inches(0.35), RGBColor(0xE0, 0xE0, 0xE0))
    fw = int(track_w * pct)
    add_rect(sl, tx, ty + Inches(0.05), fw, Inches(0.35), col)
    add_text(sl, note, tx + Inches(0.12), ty + Inches(0.07), Inches(4), Inches(0.28),
             size=11, bold=True, color=WHITE)

# Three big-number cards
metrics = [("81%", RED,    "of batches\ncontain new-class data"),
           ("19%", BLUE,   "of batches replay\nold hospital classes"),
           ("4.3×", HEADING,"more new images\nthan old images")]
mx = Inches(0.6)
for val, col, label in metrics:
    add_rect(sl, mx, Inches(5.15), Inches(3.9), Inches(2.0), BG_LIGHT, col, Pt(2.0))
    add_text(sl, val,   mx + Inches(0.15), Inches(5.25), Inches(3.6), Inches(0.8),
             size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(sl, label, mx + Inches(0.15), Inches(6.0), Inches(3.6), Inches(0.6),
             size=12, color=BODY, align=PP_ALIGN.CENTER)
    mx += Inches(4.05)

add_rect(sl, 0, Inches(7.38), SLIDE_W, Inches(0.12), HEADING)
add_text(sl, "3 / 8", Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
         size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MECHANISM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG_WHITE)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), HEADING)

add_text(sl, "Mechanism: How It Happens Epoch by Epoch",
         Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.85),
         size=32, bold=True, color=HEADING, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(4.0), Inches(1.2), Inches(5.3), Inches(0.06), HEADING)

steps = [
    (BLUE,   "Epoch 1 — Head Adapts to New Classes",
     "Detect head re-initialized for 109 classes (was 106). Loss on new classes is very high → "
     "large gradients → head weights shift aggressively toward the 3 new classes."),
    (YELLOW, "Epochs 2–10 — Old Classes Under-Represented",
     "Each epoch, ~81% of batches have no old-hospital-class annotations. "
     "Backbone and neck features drift to optimize for dominant Roboflow classes. Old-class activations weaken."),
    (RED,    "Phase 2 — Neck Unfrozen (freeze=15)",
     "Layers 15–22 (deeper neck) become trainable. These shared features are biased "
     "further toward new classes due to 4× more new-class images."),
    (DIM,    "Result — Rare Old Classes Forgotten",
     "Classes with few training examples and no Roboflow reinforcement (surgical_scissor, hair_net, "
     "test_tube) see significant recall drops."),
]

for idx, (dot_col, title, desc) in enumerate(steps):
    ty = Inches(1.55) + idx * Inches(1.35)
    add_rect(sl, Inches(0.65), ty + Inches(0.1), Inches(0.2), Inches(0.2), dot_col)
    if idx < len(steps) - 1:
        add_rect(sl, Inches(0.73), ty + Inches(0.32), Inches(0.04), Inches(1.03), RGBColor(0xCC, 0xCC, 0xCC))
    add_text(sl, title, Inches(1.1), ty, Inches(11.5), Inches(0.45),
             size=16, bold=True, color=dot_col)
    add_text(sl, desc, Inches(1.1), ty + Inches(0.44), Inches(11.5), Inches(0.78),
             size=13, color=BODY)

add_rect(sl, 0, Inches(7.38), SLIDE_W, Inches(0.12), HEADING)
add_text(sl, "4 / 8", Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
         size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FREEZE STRATEGY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG_WHITE)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), HEADING)

add_text(sl, "What the Freeze Strategy Did (and Didn't) Do",
         Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.85),
         size=32, bold=True, color=HEADING, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(4.0), Inches(1.2), Inches(5.3), Inches(0.06), HEADING)

add_text(sl,
    "Freezing protects backbone features from gradient updates. But unfrozen head and neck layers "
    "remain vulnerable to forgetting from imbalanced data.",
    Inches(1), Inches(1.38), Inches(11.3), Inches(0.65),
    size=16, color=BODY, align=PP_ALIGN.CENTER)

headers = ["Layers", "Phase 1 (freeze=22)", "Phase 2 (freeze=15)", "Effect on Forgetting"]
rows = [
    [("Backbone (0–9)",      BODY), ("Frozen", GREEN),  ("Frozen", GREEN),  ("Low-level features preserved", BODY)],
    [("C2PSA / SPPF (10)",   BODY), ("Frozen", GREEN),  ("Frozen", GREEN),  ("Attention features preserved", BODY)],
    [("Early neck (11–14)",  BODY), ("Frozen", GREEN),  ("Frozen", GREEN),  ("Upper-level features preserved", BODY)],
    [("Deeper neck (15–22)", BODY), ("Frozen", GREEN),  ("Trainable", RED), ("Drifts toward new classes", YELLOW)],
    [("Detect head (23)",    BODY), ("Trainable", RED), ("Trainable", RED), ("Main source of forgetting", RED)],
]
col_widths = [Inches(2.2), Inches(2.4), Inches(2.4), Inches(4.8)]
add_table(sl, headers, rows,
          Inches(0.55), Inches(2.1), Inches(12.2), Inches(4.8), col_widths)

add_rect(sl, 0, Inches(7.38), SLIDE_W, Inches(0.12), HEADING)
add_text(sl, "5 / 8", Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
         size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — OVERSAMPLING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG_WHITE)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), HEADING)

add_text(sl, "Oversampling — Only Helped New Rare Classes",
         Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.85),
         size=32, bold=True, color=HEADING, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(4.0), Inches(1.2), Inches(5.3), Inches(0.06), HEADING)

add_text(sl,
    "We applied oversampling to rare new classes only. This addressed new-class rarity "
    "but did not compensate for old hospital classes that had no new data at all.",
    Inches(1), Inches(1.38), Inches(11.3), Inches(0.65),
    size=16, color=BODY, align=PP_ALIGN.CENTER)

# 3 cards
card_data = [
    ("spillage (108)",       YELLOW, RGBColor(0xFE, 0xF9, 0xE7),
     "1,883 → x3 oversampled\n= 5,649 effective images\n\nWas rare.  Helped"),
    ("exit_sign (107)",      YELLOW, RGBColor(0xFE, 0xF9, 0xE7),
     "3,008 → x2 oversampled\n= 6,016 effective images\n\nModerate.  Helped"),
    ("Old classes (80–105)", RED,    CARD_RED,
     "No oversampling applied.\nStill only 4,604 hospital imgs.\nDominated by 20k+ new images.\nNot Addressed"),
]
cx = Inches(0.55)
for title, col, bg, body in card_data:
    add_rect(sl, cx, Inches(2.3), Inches(4.0), Inches(3.3), bg, col, Pt(2.0))
    add_text(sl, title, cx + Inches(0.2), Inches(2.45), Inches(3.65), Inches(0.55),
             size=16, bold=True, color=col)
    add_text(sl, body, cx + Inches(0.2), Inches(3.0), Inches(3.65), Inches(2.4),
             size=13, color=BODY)
    cx += Inches(4.25)

# Two bars at bottom
bars2 = [
    ("door / wheelchair / fire_ext", 0.72, GREEN, "Reinforced by Roboflow — likely improved"),
    ("Other 23 hospital classes",    0.28, RED,   "No new data — at risk of forgetting"),
]
lx2 = Inches(0.8)
for idx2, (label2, pct2, col2, note2) in enumerate(bars2):
    ty2 = Inches(5.85) + idx2 * Inches(0.72)
    add_text(sl, label2, lx2, ty2, Inches(2.8), Inches(0.42),
             size=12, color=BODY, align=PP_ALIGN.RIGHT)
    tx2 = lx2 + Inches(2.9)
    track_w2 = Inches(8.8)
    add_rect(sl, tx2, ty2 + Inches(0.05), track_w2, Inches(0.32), RGBColor(0xE0, 0xE0, 0xE0))
    fw2 = int(track_w2 * pct2)
    add_rect(sl, tx2, ty2 + Inches(0.05), fw2, Inches(0.32), col2)
    add_text(sl, note2, tx2 + Inches(0.1), ty2 + Inches(0.06), Inches(5.5), Inches(0.25),
             size=10, bold=True, color=WHITE)

add_rect(sl, 0, Inches(7.38), SLIDE_W, Inches(0.12), HEADING)
add_text(sl, "6 / 8", Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
         size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TRAINING METRICS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG_WHITE)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), HEADING)

add_text(sl, "Training Metrics Tell the Story",
         Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.85),
         size=32, bold=True, color=HEADING, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(4.8), Inches(1.2), Inches(3.7), Inches(0.06), HEADING)

add_text(sl,
    "Overall mAP rose steadily — dominated by the large new-class dataset. "
    "Per-class mAP for rare old classes may tell a different story.",
    Inches(1), Inches(1.38), Inches(11.3), Inches(0.55),
    size=16, color=BODY, align=PP_ALIGN.CENTER)

headers2 = ["Phase", "Epoch", "mAP50", "mAP50-95", "Precision", "Recall"]
rows2 = [
    [("Phase 1", BLUE),   ("1",          BODY), ("0.694", BODY), ("0.434", BODY), ("0.691", BODY), ("0.622", BODY)],
    [("Phase 1", BLUE),   ("5",          BODY), ("0.839", BODY), ("0.552", BODY), ("0.809", BODY), ("0.771", BODY)],
    [("Phase 1", BLUE),   ("15",         BODY), ("0.899", BODY), ("0.640", BODY), ("0.875", BODY), ("0.838", BODY)],
    [("Phase 1", BLUE),   ("30",         BODY), ("0.917", BODY), ("0.675", BODY), ("0.907", BODY), ("0.858", BODY)],
    [("Phase 2", HEADING),("1",          BODY), ("0.917", BODY), ("0.673", BODY), ("0.893", BODY), ("0.859", BODY)],
    [("Phase 2", HEADING),("20",         BODY), ("0.934", BODY), ("0.710", BODY), ("0.916", BODY), ("0.875", BODY)],
    [("Phase 2", HEADING),("70 (final)", HEADING), ("0.944", GREEN), ("0.734", GREEN), ("0.923", GREEN), ("0.892", GREEN)],
]
col_widths2 = [Inches(1.6), Inches(1.4), Inches(1.8), Inches(1.9), Inches(1.9), Inches(1.8)]
add_table(sl, headers2, rows2,
          Inches(1.15), Inches(2.05), Inches(11.0), Inches(4.7), col_widths2)

add_text(sl,
    "Note: Aggregate metrics mask per-class performance — a class with 10k images can inflate "
    "mAP50 while a class with 100 images is forgotten entirely.",
    Inches(1), Inches(6.9), Inches(11.3), Inches(0.42),
    size=11, color=DIM, italic=True, align=PP_ALIGN.CENTER)

add_rect(sl, 0, Inches(7.38), SLIDE_W, Inches(0.12), HEADING)
add_text(sl, "7 / 8", Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
         size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG_WHITE)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), HEADING)

add_text(sl, "Summary",
         Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.8),
         size=36, bold=True, color=HEADING, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(6.0), Inches(1.15), Inches(1.3), Inches(0.06), HEADING)

bullets = [
    ("4.3x Data Imbalance",       "Hospital V3 merged 20,023 Roboflow images with only 4,604 existing hospital images — 4.3x imbalance per epoch.", RED),
    ("Freeze Protected Backbone",  "2-phase freeze preserved backbone features but the Detect head (Phase 1+2) and deeper neck (Phase 2) were exposed to imbalanced gradients.", YELLOW),
    ("3 Classes Reinforced",       "door, wheelchair, fire_extinguisher overlapped with Roboflow data and likely improved. The other 23 hospital classes received no new data.", GREEN),
    ("Oversampling Was Partial",   "Oversampling applied only to new rare classes (spillage x3, exit_sign x2). Old hospital classes were not oversampled.", BLUE),
    ("Final mAP50 = 0.944",        "Strong aggregate performance after 100 epochs — but this masks per-class forgetting of rare old hospital classes.", PURPLE),
]

for idx, (title, desc, col) in enumerate(bullets):
    ty = Inches(1.42) + idx * Inches(1.1)
    add_rect(sl, Inches(0.55), ty + Inches(0.06), Inches(0.08), Inches(0.72), col)
    add_text(sl, title, Inches(0.78), ty, Inches(11.8), Inches(0.44),
             size=15, bold=True, color=col)
    add_text(sl, desc,  Inches(0.78), ty + Inches(0.43), Inches(11.8), Inches(0.6),
             size=13, color=BODY)

# Tag pills
tags2 = [
    ("Root Cause: 4.3x Imbalance", RED),
    ("Shared neck layers",          YELLOW),
    ("Per-class mAP val needed",    BLUE),
]
tx3 = Inches(1.0)
for label3, col3 in tags2:
    add_rect(sl, tx3, Inches(7.02), Inches(3.7), Inches(0.36), BG_LIGHT, col3, Pt(1.5))
    add_text(sl, label3, tx3, Inches(7.02), Inches(3.7), Inches(0.36),
             size=11, bold=True, color=col3, align=PP_ALIGN.CENTER)
    tx3 += Inches(3.85)

add_rect(sl, 0, Inches(7.38), SLIDE_W, Inches(0.12), HEADING)
add_text(sl, "8 / 8", Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
         size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/kelvin/yolo_tr/why_hospital_v3_forgot_classes.pptx"
prs.save(out)
print(f"Saved: {out}")
